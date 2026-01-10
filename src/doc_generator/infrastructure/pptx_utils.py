"""
PowerPoint generation utilities using python-pptx.

Provides helper functions for creating PPTX presentations programmatically.
Optimized for corporate/executive presentations.
"""

from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Corporate-ready color theme - Professional and polished
THEME_COLORS = {
    "ink": RGBColor(26, 26, 46),          # Deep navy for text
    "muted": RGBColor(74, 85, 104),       # Sophisticated gray
    "accent": RGBColor(37, 99, 235),      # Professional blue
    "accent_dark": RGBColor(29, 78, 216), # Darker blue
    "teal": RGBColor(13, 148, 136),       # Modern teal
    "success": RGBColor(16, 185, 129),    # Green for positive
    "warning": RGBColor(245, 158, 11),    # Amber for warnings
    "background": RGBColor(255, 255, 255),
    "light_bg": RGBColor(248, 250, 252),  # Clean white background
    "dark_bg": RGBColor(30, 41, 59),      # Slate background for contrast
    "gradient_start": RGBColor(37, 99, 235),  # Blue gradient start
    "gradient_end": RGBColor(13, 148, 136),   # Teal gradient end
}


def create_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation with 16:9 layout.

    Returns:
        Presentation object
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    logger.debug("Created new PowerPoint presentation (16:9)")
    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """
    Add corporate-style title slide to presentation.

    Features a clean, professional design with modern typography.

    Args:
        prs: Presentation object
        title: Title text
        subtitle: Subtitle text (optional)
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(slide_layout)

    # Clean white background for corporate feel
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME_COLORS["background"]

    # Accent bar at top (thin for professional look)
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.12)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    accent_bar.line.fill.background()

    # Left accent bar for visual interest
    left_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.4), Inches(1.3),
        Inches(0.08), Inches(2.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    left_bar.line.fill.background()

    # Title - large, professional, left-aligned
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.4),
        Inches(8.5), Inches(2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = THEME_COLORS["ink"]
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.4),
            Inches(8.5), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = THEME_COLORS["muted"]
        p.alignment = PP_ALIGN.LEFT

    # Bottom accent line (left-aligned for corporate look)
    bottom_line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.7), Inches(4.6),
        Inches(3), Inches(0.04)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = THEME_COLORS["teal"]
    bottom_line.line.fill.background()

    logger.debug(f"Added corporate title slide: {title}")


def add_content_slide(
    prs: Presentation,
    title: str,
    content: list[str],
    is_bullets: bool = True,
    speaker_notes: str = ""
) -> None:
    """
    Add corporate-style content slide with title and bullet points.

    Features clean design with professional typography and spacing.

    Args:
        prs: Presentation object
        title: Slide title
        content: List of content items
        is_bullets: Whether to format as bullets (default True)
        speaker_notes: Optional speaker notes for the slide
    """
    slide_layout = prs.slide_layouts[6]  # Blank for custom styling
    slide = prs.slides.add_slide(slide_layout)

    # Add thin accent bar at top
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.06)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    accent_bar.line.fill.background()

    # Title with left-aligned corporate style
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = THEME_COLORS["ink"]
    p.font.bold = True

    # Underline for title
    title_line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.1),
        Inches(1.5), Inches(0.03)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = THEME_COLORS["accent"]
    title_line.line.fill.background()

    # Content area (corporate-style with professional spacing)
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(9), Inches(3.9)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Clean bullet markers if present
        clean_item = item.lstrip("•-* ").strip()

        if is_bullets:
            p.text = f"•  {clean_item}"
            p.font.size = Pt(20)
            p.font.color.rgb = THEME_COLORS["ink"]
            p.space_after = Pt(14)
        else:
            p.text = clean_item
            p.font.size = Pt(18)
            p.font.color.rgb = THEME_COLORS["muted"]
            p.space_after = Pt(12)

    # Add speaker notes if provided
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = speaker_notes

    logger.debug(f"Added corporate content slide: {title} ({len(content)} items)")


def add_section_header_slide(prs: Presentation, section_title: str) -> None:
    """
    Add corporate-style section header slide with professional dark background.

    Args:
        prs: Presentation object
        section_title: Section heading text
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Professional dark slate background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME_COLORS["dark_bg"]

    # Accent bar on left side
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(0.15), prs.slide_height
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    accent_bar.line.fill.background()

    # Section indicator line
    indicator = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(2.2),
        Inches(2), Inches(0.04)
    )
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = THEME_COLORS["accent"]
    indicator.line.fill.background()

    # Section title - large and bold
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.4),
        Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(42)
    p.font.color.rgb = THEME_COLORS["background"]
    p.font.bold = True

    # Bottom decorative line
    bottom_line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(4.2),
        Inches(1.5), Inches(0.03)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = THEME_COLORS["teal"]
    bottom_line.line.fill.background()

    logger.debug(f"Added corporate section header: {section_title}")


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str = ""
) -> None:
    """
    Add slide with image and optional caption.

    Args:
        prs: Presentation object
        title: Slide title
        image_path: Path to image file
        caption: Image caption (optional)
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLORS["ink"]
    p.font.bold = True

    # Add image (centered)
    if image_path.exists():
        left = Inches(1.5)
        top = Inches(1.5)
        max_width = Inches(7)
        Inches(3.5)

        slide.shapes.add_picture(
            str(image_path),
            left, top,
            width=max_width
        )

    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(0.5)

        caption_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = caption_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = THEME_COLORS["muted"]
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    logger.debug(f"Added image slide: {title}")


def add_executive_summary_slide(
    prs: Presentation,
    title: str,
    summary_points: list[str]
) -> None:
    """
    Add executive summary slide with key takeaways.

    Args:
        prs: Presentation object
        title: Slide title (usually "Executive Summary")
        summary_points: List of key summary points
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Clean white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME_COLORS["background"]

    # Accent bar at top
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.06)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLORS["ink"]
    p.font.bold = True

    # Title underline
    title_line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(0.9),
        Inches(1.8), Inches(0.03)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = THEME_COLORS["accent"]
    title_line.line.fill.background()

    # Key points with numbered indicators
    for i, point in enumerate(summary_points[:5]):  # Max 5 points
        y_pos = 1.1 + i * 0.8

        # Number circle indicator
        num_circle = slide.shapes.add_shape(
            9,  # Oval
            Inches(0.5), Inches(y_pos),
            Inches(0.35), Inches(0.35)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = THEME_COLORS["accent"]
        num_circle.line.fill.background()

        # Number text
        num_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos + 0.02),
            Inches(0.35), Inches(0.35)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i + 1}"
        p.font.size = Pt(16)
        p.font.color.rgb = THEME_COLORS["background"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Point text
        point_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(y_pos),
            Inches(8.4), Inches(0.7)
        )
        tf = point_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        clean_point = point.lstrip("•-* ").strip()
        p.text = clean_point
        p.font.size = Pt(16)
        p.font.color.rgb = THEME_COLORS["ink"]

    logger.debug(f"Added executive summary slide: {len(summary_points)} points")


def add_chart_slide(
    prs: Presentation,
    title: str,
    svg_path: Path,
    caption: str = ""
) -> None:
    """
    Add slide with SVG chart (converted to PNG).

    Args:
        prs: Presentation object
        title: Slide title
        svg_path: Path to SVG file
        caption: Optional caption
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3),
        Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLORS["ink"]
    p.font.bold = True

    # Convert SVG to PNG and add to slide
    if svg_path.exists():
        try:
            import tempfile

            import cairosvg

            # Convert SVG to PNG
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                png_path = tmp.name
                cairosvg.svg2png(url=str(svg_path), write_to=png_path, scale=2.0)

                slide.shapes.add_picture(
                    png_path,
                    Inches(1), Inches(1.2),
                    width=Inches(8)
                )

                # Cleanup temp file
                Path(png_path).unlink(missing_ok=True)
        except Exception as e:
            logger.warning(f"Failed to add chart: {e}")

    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(5),
            Inches(8.8), Inches(0.4)
        )
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = THEME_COLORS["muted"]
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    logger.debug(f"Added chart slide: {title}")


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_content: list[str],
    right_content: list[str],
    left_title: str = "",
    right_title: str = ""
) -> None:
    """
    Add two-column comparison slide.

    Args:
        prs: Presentation object
        title: Main slide title
        left_content: Left column bullet points
        right_content: Right column bullet points
        left_title: Left column header
        right_title: Right column header
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = THEME_COLORS["accent"]
    accent_bar.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3),
        Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLORS["ink"]
    p.font.bold = True

    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.1),
            Inches(4.2), Inches(0.5)
        )
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(18)
        p.font.color.rgb = THEME_COLORS["accent"]
        p.font.bold = True

    # Left column content
    left_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.6),
        Inches(4.2), Inches(3.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content[:6]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item.lstrip('•-* ').strip()}"
        p.font.size = Pt(16)
        p.font.color.rgb = THEME_COLORS["ink"]
        p.space_after = Pt(8)

    # Vertical divider
    divider = slide.shapes.add_shape(
        1, Inches(4.9), Inches(1.1),
        Inches(0.02), Inches(4)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = THEME_COLORS["muted"]
    divider.line.fill.background()

    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.1),
            Inches(4.2), Inches(0.5)
        )
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(18)
        p.font.color.rgb = THEME_COLORS["teal"]
        p.font.bold = True

    # Right column content
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.6),
        Inches(4.2), Inches(3.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content[:6]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item.lstrip('•-* ').strip()}"
        p.font.size = Pt(16)
        p.font.color.rgb = THEME_COLORS["ink"]
        p.space_after = Pt(8)

    logger.debug(f"Added two-column slide: {title}")


def save_presentation(prs: Presentation, output_path: Path) -> None:
    """
    Save presentation to file.

    Args:
        prs: Presentation object
        output_path: Path to output file

    Raises:
        IOError: If save fails
    """
    try:
        # Ensure parent directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs.save(str(output_path))
        logger.info(f"Saved presentation: {output_path}")

    except Exception as e:
        logger.error(f"Failed to save presentation: {e}")
        raise IOError(f"Failed to save presentation: {e}")
