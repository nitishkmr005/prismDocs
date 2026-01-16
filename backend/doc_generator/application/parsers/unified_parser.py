"""
Unified parser using Docling for multiple document formats.

Handles PDF, DOCX, PPTX, XLSX, and images with advanced OCR support.
Falls back to lighter parsers (pypdf, python-docx) when Docling is unavailable.
"""

import os
from pathlib import Path
from typing import Tuple

from loguru import logger

from ...domain.exceptions import ParseError
from ...infrastructure.parsers.docling import (
    convert_document_to_markdown,
    is_docling_available,
)


class UnifiedParser:
    """
    Parser for PDF, DOCX, PPTX, images using Docling.

    Uses IBM Research's Docling library for advanced document parsing
    with OCR, table extraction, and layout analysis.

    Falls back to lighter parsers when Docling is not available:
    - PDF: pypdf
    - DOCX: python-docx
    - PPTX: python-pptx
    """

    def __init__(self):
        """Initialize parser and check Docling availability."""
        self._docling_available = is_docling_available()
        self._docling_enabled = os.getenv("DOCLING_ENABLED", "true").lower() == "true"

        if not self._docling_available:
            logger.warning(
                "Docling not available - will use fallback parsers (pypdf, python-docx)"
            )
        elif not self._docling_enabled:
            logger.info(
                "Docling disabled via DOCLING_ENABLED=false - using fallback parsers"
            )

    def parse(self, input_path: str | Path) -> Tuple[str, dict]:
        """
        Parse document using Docling or fallback parsers.

        Args:
            input_path: Path to input file

        Returns:
            Tuple of (markdown_content, metadata)

        Raises:
            ParseError: If parsing fails
        """
        path = Path(input_path)

        if not path.exists():
            raise ParseError(f"File not found: {path}")

        suffix = path.suffix.lower()

        # Use Docling if available and enabled
        if self._docling_available and self._docling_enabled:
            return self._parse_with_docling(path)

        # Fallback parsers based on file type
        logger.info(f"Using fallback parser for: {path.name}")

        if suffix == ".pdf":
            return self._parse_pdf_fallback(path)
        elif suffix == ".docx":
            return self._parse_docx_fallback(path)
        elif suffix == ".pptx":
            return self._parse_pptx_fallback(path)
        elif suffix in [".png", ".jpg", ".jpeg", ".tiff"]:
            return self._parse_image_fallback(path)
        elif suffix == ".xlsx":
            return self._parse_xlsx_fallback(path)
        else:
            raise ParseError(
                f"Unsupported format '{suffix}' and Docling is not available. "
                "Install Docling for advanced document parsing: pip install docling"
            )

    def _parse_with_docling(self, path: Path) -> Tuple[str, dict]:
        """Parse using Docling (full functionality)."""
        logger.info(f"Parsing with Docling: {path.name}")
        try:
            content, metadata = convert_document_to_markdown(path)
            logger.info(
                f"Docling parsing completed: {len(content)} chars, "
                f"{metadata.get('pages', 'N/A')} pages"
            )
            return content, metadata
        except Exception as e:
            logger.error(f"Docling parsing failed for {path}: {e}")
            raise ParseError(f"Failed to parse document with Docling: {e}")

    def _parse_pdf_fallback(self, path: Path) -> Tuple[str, dict]:
        """Parse PDF using pypdf (basic text extraction, no OCR)."""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ParseError("pypdf not installed. Run: pip install pypdf")

        logger.info(f"Parsing PDF with pypdf (fallback): {path.name}")

        try:
            reader = PdfReader(str(path))
            pages_text = []

            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages_text.append(f"<!-- Page {i+1} -->\n{text}")

            content = "\n\n".join(pages_text)

            metadata = {
                "title": path.stem,
                "pages": len(reader.pages),
                "source_file": str(path),
                "parser": "pypdf-fallback",
            }

            logger.info(
                f"pypdf parsing completed: {len(content)} chars, {len(reader.pages)} pages"
            )
            return content, metadata

        except Exception as e:
            logger.error(f"pypdf parsing failed for {path}: {e}")
            raise ParseError(f"Failed to parse PDF: {e}")

    def _parse_docx_fallback(self, path: Path) -> Tuple[str, dict]:
        """Parse DOCX using python-docx."""
        try:
            from docx import Document
        except ImportError:
            raise ParseError("python-docx not installed. Run: pip install python-docx")

        logger.info(f"Parsing DOCX with python-docx (fallback): {path.name}")

        try:
            doc = Document(str(path))
            paragraphs = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Check for heading styles
                    if para.style and para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        try:
                            level_num = int(level)
                            heading_prefix = "#" * min(level_num, 6)
                            paragraphs.append(f"{heading_prefix} {text}")
                        except ValueError:
                            paragraphs.append(f"## {text}")
                    else:
                        paragraphs.append(text)

            # Extract tables
            for table in doc.tables:
                table_md = self._table_to_markdown(table)
                if table_md:
                    paragraphs.append(table_md)

            content = "\n\n".join(paragraphs)

            metadata = {
                "title": path.stem,
                "source_file": str(path),
                "parser": "python-docx-fallback",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
            }

            logger.info(f"python-docx parsing completed: {len(content)} chars")
            return content, metadata

        except Exception as e:
            logger.error(f"python-docx parsing failed for {path}: {e}")
            raise ParseError(f"Failed to parse DOCX: {e}")

    def _parse_pptx_fallback(self, path: Path) -> Tuple[str, dict]:
        """Parse PPTX using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ParseError("python-pptx not installed. Run: pip install python-pptx")

        logger.info(f"Parsing PPTX with python-pptx (fallback): {path.name}")

        try:
            prs = Presentation(str(path))
            slides_content = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"## Slide {i}"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                slides_content.append("\n\n".join(slide_text))

            content = "\n\n---\n\n".join(slides_content)

            metadata = {
                "title": path.stem,
                "source_file": str(path),
                "parser": "python-pptx-fallback",
                "slides": len(prs.slides),
            }

            logger.info(
                f"python-pptx parsing completed: {len(content)} chars, {len(prs.slides)} slides"
            )
            return content, metadata

        except Exception as e:
            logger.error(f"python-pptx parsing failed for {path}: {e}")
            raise ParseError(f"Failed to parse PPTX: {e}")

    def _parse_image_fallback(self, path: Path) -> Tuple[str, dict]:
        """Handle images without OCR (returns placeholder)."""
        logger.warning(
            f"Image parsing without Docling - no OCR available for: {path.name}"
        )

        content = f"![Image: {path.name}]({path.name})\n\n*Note: OCR not available. Install Docling for text extraction from images.*"

        metadata = {
            "title": path.stem,
            "source_file": str(path),
            "parser": "image-placeholder",
            "warning": "No OCR - Docling not available",
        }

        return content, metadata

    def _parse_xlsx_fallback(self, path: Path) -> Tuple[str, dict]:
        """Handle XLSX without Docling (limited support)."""
        logger.warning(f"XLSX parsing without Docling is limited: {path.name}")

        content = f"*Excel file: {path.name}*\n\n*Note: For full XLSX support, install Docling.*"

        metadata = {
            "title": path.stem,
            "source_file": str(path),
            "parser": "xlsx-placeholder",
            "warning": "Limited XLSX support - Docling not available",
        }

        return content, metadata

    def _table_to_markdown(self, table) -> str:
        """Convert a python-docx table to markdown format."""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")

            if len(rows) >= 1:
                # Add header separator after first row
                num_cols = len(table.rows[0].cells)
                separator = "| " + " | ".join(["---"] * num_cols) + " |"
                rows.insert(1, separator)

            return "\n".join(rows)
        except Exception:
            return ""
